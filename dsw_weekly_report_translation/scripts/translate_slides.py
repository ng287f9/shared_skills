#!/usr/bin/env python3
"""Translate a Chinese DSW weekly-report .pptx into an English-only version.

Usage:
    python translate_slides.py --input <src.pptx> --mapping <mapping.json> \
        [--font Arial] [--out <out.pptx>]

mapping.json is a flat dict:
    {"<exact original paragraph full run-text>": "<English translation>", ...}
The value may also be an object {"en": "...", "sz": 13.5, "spcBef": 4}:
  - "sz": force the run font size (points) — use only when the English text
    needs a small size reduction to fit its shape.
  - "spcBef": set the paragraph's space-before to this many points — use to
    reclaim a few points per bullet when a dense bullet slide runs long (the
    English text is usually a little taller than the Chinese).
Plain strings keep the original formatting.

Only paragraphs whose full concatenated run-text exactly matches a key are
changed. Everything else (images/screenshots, layout, colors, paragraph props,
empty paragraphs, IDs, numbers) is preserved exactly.

Font handling:
  * every translated run gets --font (default Arial) as BOTH latin and EA
    typeface, and copies the original run's size/bold/color;
  * small fixed-width label shapes (< 3.5 in wide) get a single-line fit pass:
    if the English text is wider than the shape, the font size is reduced
    until it fits (min 6.5 pt) so nothing overflows.

Output: <input_stem>_en.pptx unless --out is given.
"""
import argparse
import copy
import json
import sys
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_PT = 12700
SMALL_SHAPE_WIDTH_EMU = 3.5 * 914400  # shapes narrower than this are single-line labels
FIT_MIN_PT = 6.5
FIT_STEP_PT = 0.5
XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

ARIAL_PATHS = {
    False: "/System/Library/Fonts/Supplemental/Arial.ttf",
    True: "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
}

_font_cache = {}


def measure_width(text: str, size_pt: float, bold: bool) -> float:
    """Approximate text width in points at the given size (rendered 4x)."""
    key = (round(size_pt, 2), bold)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(ARIAL_PATHS[bold], int(round(size_pt * 4)))
    return _font_cache[key].getlength(text) / 4.0


def _set_fonts(rpr, font: str) -> None:
    """Force latin+ea typefaces on an a:rPr element."""
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is not None:
            rpr.remove(el)
    latin = rpr.makeelement(qn("a:latin"), {})
    latin.set("typeface", font)
    rpr.append(latin)
    ea = rpr.makeelement(qn("a:ea"), {})
    ea.set("typeface", font)
    rpr.append(ea)


def _shape_insets(shape) -> tuple[int, int]:
    """left/right text insets in EMU from the bodyPr, or DrawingML defaults."""
    try:
        txbody = shape.text_frame._txBody
    except Exception:
        return 91440, 91440
    bodypr = txbody.find(qn("a:bodyPr"))
    if bodypr is None:
        return 91440, 91440
    l = bodypr.get("lIns")
    r = bodypr.get("rIns")
    l_ins = int(l) if l is not None else 91440
    r_ins = int(r) if r is not None else 91440
    return l_ins, r_ins


def _fit_size(shape, text: str, size_pt: float, bold: bool, single_para: bool) -> float:
    """Reduce size_pt so text fits one line in a single-line label shape.

    Applies to small fixed-width label shapes (< 3.5 in) and to any shape that
    holds exactly one paragraph (e.g. a cover headline) whose single-line text
    would otherwise wrap or overflow. Body bullets (multi-paragraph, wrapping)
    are never shrunk.
    """
    if shape is None or shape.width is None:
        return size_pt
    is_small = shape.width < SMALL_SHAPE_WIDTH_EMU
    if not is_small and not single_para:
        return size_pt
    l_ins, r_ins = _shape_insets(shape)
    avail_pt = (shape.width - l_ins - r_ins) / EMU_PER_PT
    if avail_pt <= 0:
        return size_pt
    while size_pt - FIT_STEP_PT >= FIT_MIN_PT and measure_width(text, size_pt, bold) > avail_pt:
        size_pt = round(size_pt - FIT_STEP_PT, 1)
    return size_pt


def _set_space_before(p_el, pts: float) -> None:
    """Set the paragraph's space-before (<a:spcBef>) to a value in points."""
    ppr = p_el.find(qn("a:pPr"))
    if ppr is None:
        ppr = p_el.makeelement(qn("a:pPr"), {})
        p_el.insert(0, ppr)
    spcBef = ppr.find(qn("a:spcBef"))
    if spcBef is None:
        spcBef = ppr.makeelement(qn("a:spcBef"), {})
        lnspc = ppr.find(qn("a:lnSpc"))
        if lnspc is not None:
            lnspc.addnext(spcBef)
        else:
            ppr.insert(0, spcBef)
    for child in list(spcBef):
        spcBef.remove(child)
    spcBef.append(spcBef.makeelement(qn("a:spcPts"), {"val": str(int(round(pts * 100)))}))


def _new_run_xml(p_el, text: str, font: str, old_runs, shape, single_para: bool,
                 override_sz=None, spc_bef_pts=None) -> None:
    """Replace all runs in paragraph <a:p> with one translated run."""
    # --- copy first run's rPr (keeps size/bold/color/inheritance) ---
    first_rpr = None
    if old_runs:
        r0 = old_runs[0]
        rpr_el = r0.find(qn("a:rPr"))
        if rpr_el is not None:
            first_rpr = copy.deepcopy(rpr_el)

    old_sz = None
    bold = False
    if first_rpr is not None:
        sz_attr = first_rpr.get("sz")
        old_sz = int(sz_attr) / 100.0 if sz_attr is not None else None
        bold = first_rpr.get("b") in ("1", "true")

    # --- fit pass (single-line labels only); explicit override wins ---
    size_pt = old_sz
    if size_pt is not None:
        if override_sz is not None:
            size_pt = float(override_sz)
        else:
            size_pt = _fit_size(shape, text, size_pt, bold, single_para)

    # --- strip old runs / breaks / fields ---
    for child in list(p_el):
        if child.tag in (qn("a:r"), qn("a:br"), qn("a:fld")):
            p_el.remove(child)

    # --- insert new run after pPr (or at start) ---
    r = p_el.makeelement(qn("a:r"), {})
    if first_rpr is not None:
        if size_pt is not None and old_sz is not None and size_pt != old_sz:
            first_rpr.set("sz", str(int(round(size_pt * 100))))
        _set_fonts(first_rpr, font)
        r.append(first_rpr)
    else:
        rpr = p_el.makeelement(qn("a:rPr"), {})
        if size_pt is not None:
            rpr.set("sz", str(int(round(size_pt * 100))))
        _set_fonts(rpr, font)
        r.append(rpr)
    t = p_el.makeelement(qn("a:t"), {})
    t.text = text
    if text != text.strip():
        t.set(XML_SPACE, "preserve")
    r.append(t)

    ppr = p_el.find(qn("a:pPr"))
    if ppr is not None:
        ppr.addnext(r)
    else:
        p_el.insert(0, r)

    if spc_bef_pts is not None:
        _set_space_before(p_el, spc_bef_pts)


def _iter_text_frames(shapes):
    for shp in shapes:
        st = shp.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_text_frames(shp.shapes)
            except Exception:
                pass
        elif st == MSO_SHAPE_TYPE.TABLE:
            try:
                for row in shp.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame, shp
            except Exception:
                pass
        elif shp.has_text_frame:
            yield shp.text_frame, shp


def translate(pptx_path: str, mapping: dict, font: str, out_path: str) -> int:
    prs = Presentation(pptx_path)
    replaced = 0
    misses = []
    for slide in prs.slides:
        for tf, shape in _iter_text_frames(slide.shapes):
            txbody = tf._txBody
            paras = txbody.findall(qn("a:p"))
            non_empty = [
                p for p in paras
                if "".join(t.text or "" for t in p.iter(qn("a:t"))).strip()
            ]
            single_para = len(non_empty) == 1
            for p_el in paras:
                full = "".join(t.text or "" for t in p_el.iter(qn("a:t")))
                if full in mapping:
                    old_runs = p_el.findall(qn("a:r"))
                    value = mapping[full]
                    if isinstance(value, dict):
                        text = value["en"]
                        override_sz = value.get("sz")
                        spc_bef_pts = value.get("spcBef")
                    else:
                        text, override_sz, spc_bef_pts = value, None, None
                    _new_run_xml(p_el, text, font, old_runs, shape, single_para,
                                 override_sz, spc_bef_pts)
                    replaced += 1
    prs.save(out_path)
    return replaced


def main() -> int:
    ap = argparse.ArgumentParser(description="Translate a Chinese DSW weekly report to English-only.")
    ap.add_argument("--input", required=True)
    ap.add_argument("--mapping", required=True)
    ap.add_argument("--font", default="Arial")
    ap.add_argument("--out", default=None)
    args = ap.parse_args()

    src = Path(args.input)
    if not src.exists():
        sys.exit(f"input not found: {src}")
    mapping = json.loads(Path(args.mapping).read_text(encoding="utf-8"))
    out = Path(args.out) if args.out else src.with_name(src.stem + "_en.pptx")

    n = translate(str(src), mapping, args.font, str(out))
    print(f"replaced {n} paragraphs -> {out}")


if __name__ == "__main__":
    main()
