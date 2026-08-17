#!/usr/bin/env python
"""
Phase 3: Insert template Flight Test Program Progress slide as a new slide
after the rebuilt data slide in the output. Shifts existing tables slides
(slides 6-7 → 7-8) and inserts the template's slide 5 as slide 6.

Usage:
    python insert_template_slide.py <input.pptx> <output.pptx> [<template.pptx>]
"""

import sys, os, shutil, re, zipfile
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation as PptxPresentation

def safe_move(src, dst):
    if os.path.exists(dst): os.remove(dst)
    if os.path.exists(src): os.rename(src, dst)

def safe_copy(src, dst):
    if os.path.exists(dst): os.remove(dst)
    if os.path.exists(src): shutil.copy(src, dst)

def file_text(p):
    return open(p, encoding='utf-8').read() if os.path.exists(p) else ''

def write_file(p, s):
    open(p, 'w', encoding='utf-8').write(s)

SELF_DIR = os.path.dirname(os.path.abspath(__file__))
SKILL_DIR = os.path.dirname(SELF_DIR)
DEFAULT_TEMPLATE = os.path.join(SKILL_DIR, 'references', 'template.pptx')

def process(input_pptx, output_pptx, template_pptx=None):
    if template_pptx is None:
        template_pptx = DEFAULT_TEMPLATE
    if not os.path.exists(template_pptx):
        print(f'ERROR: Template not found at {template_pptx}')
        sys.exit(1)

    work = os.path.join(os.path.dirname(output_pptx), '_insert_work')
    if os.path.exists(work):
        shutil.rmtree(work)

    tmpl = os.path.join(os.path.dirname(output_pptx), '_tmpl_tmp')
    if os.path.exists(tmpl):
        shutil.rmtree(tmpl)

    # Unpack both
    with zipfile.ZipFile(input_pptx, 'r') as z:
        z.extractall(work)
    with zipfile.ZipFile(template_pptx, 'r') as z:
        z.extractall(tmpl)

    SLIDES = os.path.join(work, 'ppt', 'slides')
    RELS = os.path.join(SLIDES, '_rels')
    NOTES = os.path.join(work, 'ppt', 'notesSlides')
    NRELS = os.path.join(NOTES, '_rels')
    MEDIA = os.path.join(work, 'ppt', 'media')
    T_SLIDES = os.path.join(tmpl, 'ppt', 'slides')
    T_RELS = os.path.join(T_SLIDES, '_rels')
    T_NOTES = os.path.join(tmpl, 'ppt', 'notesSlides')
    T_MEDIA = os.path.join(tmpl, 'ppt', 'media')

    # Shift slides 7→8, 6→7
    print('Shifting slides 7→8, 6→7...')
    for old_n, new_n in [(7, 8), (6, 7)]:
        for ext in ['.xml', '.xml.rels']:
            target = os.path.join(SLIDES if ext == '.xml' else RELS, f'slide{new_n}{ext}')
            if os.path.exists(target): os.remove(target)
        target_note = os.path.join(NOTES, f'notesSlide{new_n}.xml')
        if os.path.exists(target_note): os.remove(target_note)
        target_nr = os.path.join(NRELS, f'notesSlide{new_n}.xml.rels')
        if os.path.exists(target_nr): os.remove(target_nr)

        safe_move(os.path.join(SLIDES, f'slide{old_n}.xml'), os.path.join(SLIDES, f'slide{new_n}.xml'))
        safe_move(os.path.join(RELS, f'slide{old_n}.xml.rels'), os.path.join(RELS, f'slide{new_n}.xml.rels'))
        safe_move(os.path.join(NOTES, f'notesSlide{old_n}.xml'), os.path.join(NOTES, f'notesSlide{new_n}.xml'))
        safe_move(os.path.join(NRELS, f'notesSlide{old_n}.xml.rels'), os.path.join(NRELS, f'notesSlide{new_n}.xml.rels'))

        rpath = os.path.join(RELS, f'slide{new_n}.xml.rels')
        if os.path.exists(rpath):
            d = file_text(rpath)
            d = d.replace(f'notesSlide{old_n}', f'notesSlide{new_n}')
            write_file(rpath, d)

    # Copy template slide 5 → slide 6
    print('Copying template slide 5 → slide 6...')
    safe_copy(os.path.join(T_SLIDES, 'slide5.xml'), os.path.join(SLIDES, 'slide6.xml'))
    safe_copy(os.path.join(T_RELS, 'slide5.xml.rels'), os.path.join(RELS, 'slide6.xml.rels'))
    safe_copy(os.path.join(T_NOTES, 'notesSlide5.xml'), os.path.join(NOTES, 'notesSlide6.xml'))
    safe_copy(os.path.join(T_NOTES, '_rels', 'notesSlide5.xml.rels'), os.path.join(NRELS, 'notesSlide6.xml.rels'))

    d = file_text(os.path.join(RELS, 'slide6.xml.rels'))
    d = d.replace('notesSlide5', 'notesSlide6')
    write_file(os.path.join(RELS, 'slide6.xml.rels'), d)

    for m in re.findall(r'Target="\.\./media/(.+?)"', d):
        src = os.path.join(T_MEDIA, m)
        if os.path.exists(src) and not os.path.exists(os.path.join(MEDIA, m)):
            safe_copy(src, os.path.join(MEDIA, m))

    # Update [Content_Types].xml
    print('Updating Content_Types...')
    ct = file_text(os.path.join(work, '[Content_Types].xml'))
    for fname, ctype in [
        ('slide8.xml', 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'),
        ('notesSlide6.xml', 'application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml'),
        ('notesSlide8.xml', 'application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml'),
    ]:
        if fname not in ct:
            if fname.startswith('slide'):
                part_path = '/ppt/slides/' + fname
            elif fname.startswith('notesSlide'):
                part_path = '/ppt/notesSlides/' + fname
            else:
                part_path = '/ppt/' + fname
            ct = ct.replace('</Types>', f'<Override PartName="{part_path}" ContentType="{ctype}"/>\n</Types>')
    write_file(os.path.join(work, '[Content_Types].xml'), ct)

    # Update presentation.xml.rels
    print('Updating presentation rels...')
    pr = os.path.join(work, 'ppt', '_rels', 'presentation.xml.rels')
    tree = etree.parse(pr)
    root = tree.getroot()
    R = 'http://schemas.openxmlformats.org/package/2006/relationships'
    for rel in root.findall(f'{{{R}}}Relationship'):
        t = rel.get('Target', '')
        if t == 'slides/slide6.xml':
            rel.set('Target', 'slides/slide7.xml')
        elif t == 'slides/slide7.xml':
            rel.set('Target', 'slides/slide8.xml')
    mx = 0
    for rel in root.findall(f'{{{R}}}Relationship'):
        rid = rel.get('Id', '')
        m = re.search(r'rId(\d+)', rid)
        if m: mx = max(mx, int(m.group(1)))
    new_rid = f'rId{mx + 1}'
    new_rel = etree.SubElement(root, f'{{{R}}}Relationship')
    new_rel.set('Id', new_rid)
    new_rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
    new_rel.set('Target', 'slides/slide6.xml')
    tree.write(pr, xml_declaration=True, encoding='UTF-8', standalone=True)
    print(f'  Added {new_rid} → slides/slide6.xml')

    # Update presentation.xml sldIdLst
    print('Updating sldIdLst...')
    px = os.path.join(work, 'ppt', 'presentation.xml')
    tree2 = etree.parse(px)
    root2 = tree2.getroot()
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    RD_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    # Register namespaces for proper serialization
    etree.register_namespace('p', P_NS)
    etree.register_namespace('r', RD_NS)
    etree.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')

    sldIdLst = root2.find(f'{{{P_NS}}}sldIdLst')
    for s in list(sldIdLst):
        sldIdLst.remove(s)

    # Build rId mapping: slide1-5=rId5-9, slide6=new, slide7=rId10, slide8=rId11
    rids = ['rId5', 'rId6', 'rId7', 'rId8', 'rId9', new_rid, 'rId10', 'rId11']
    for i, rid in enumerate(rids):
        sldId = etree.SubElement(sldIdLst, f'{{{P_NS}}}sldId')
        sldId.set('id', str(256 + i))
        sldId.set(f'{{{RD_NS}}}id', rid)

    tree2.write(px, xml_declaration=True, encoding='UTF-8', standalone=True)

    # Repack
    print('Repacking...')
    if os.path.exists(output_pptx): os.remove(output_pptx)
    with zipfile.ZipFile(output_pptx, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root_dir, dirs, files in os.walk(work):
            for f in files:
                zf.write(os.path.join(root_dir, f), os.path.relpath(os.path.join(root_dir, f), work))

    # Cleanup
    shutil.rmtree(work)
    shutil.rmtree(tmpl)

    # Verify
    prs = PptxPresentation(output_pptx)
    print(f'\nOutput: {len(prs.slides)} slides')
    for i, slide in enumerate(prs.slides, 1):
        first = ''
        for shape in slide.shapes:
            if shape.has_text_frame:
                first = shape.text_frame.text.strip()[:80]
                break
            if shape.has_table:
                first = '[TABLE] ' + shape.table.cell(0,0).text.replace('\n',' | ')[:60]
                break
        tbl = sum(1 for s in slide.shapes if s.has_table)
        img = sum(1 for s in slide.shapes if s.shape_type == 13)
        print(f'  Slide {i}: {tbl}T {img}I → {first}')
    print(f'\nSize: {os.path.getsize(output_pptx)} bytes')

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: python insert_template_slide.py <input.pptx> <output.pptx> [<template.pptx>]')
        sys.exit(1)
    tpl = sys.argv[3] if len(sys.argv) > 3 else None
    process(sys.argv[1], sys.argv[2], tpl)
