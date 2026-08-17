#!/usr/bin/env python
"""
Phase 2: Rebuild Flight Test Program Progress (slide 5) screenshot tables
as native bilingual tables. Uses template slide 5 for formatting reference.
Reads CW-week source data from the original screenshots via OCR.

Usage:
    python rebuild_flight_test_progress.py <input.pptx> <output.pptx>
"""

import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

def hex_to_rgb(h):
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_borders(cell, w=12700, color='000000'):
    tcPr = cell._tc.get_or_add_tcPr()
    for side in ['lnL', 'lnR', 'lnT', 'lnB']:
        for old in tcPr.findall(qn(f'a:{side}')):
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(f'a:{side}'))
        ln.set('w', str(w)); ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', color)
        etree.SubElement(ln, qn('a:prstDash')).set('val', 'solid')
        etree.SubElement(ln, qn('a:round'))
        for end_n in ['headEnd', 'tailEnd']:
            e = etree.SubElement(ln, qn(f'a:{end_n}'))
            e.set('type', 'none'); e.set('w', 'med'); e.set('len', 'med')

def set_cell(cell, en_text, cn_text, font_size=Pt(7), bold=False,
             color='000000', fill=None, align=PP_ALIGN.CENTER):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p_en = tf.paragraphs[0]
    p_en.alignment = align
    p_en.space_before = Pt(0); p_en.space_after = Pt(1)
    run = p_en.add_run()
    run.text = en_text
    run.font.size = font_size; run.font.bold = bold
    run.font.color.rgb = RGBColor(*hex_to_rgb(color))
    run.font.name = 'Arial'
    if cn_text:
        p_cn = tf.add_paragraph()
        p_cn.alignment = align
        p_cn.space_before = Pt(0); p_cn.space_after = Pt(0)
        run2 = p_cn.add_run()
        run2.text = cn_text
        run2.font.size = font_size; run2.font.bold = bold
        run2.font.color.rgb = RGBColor(*hex_to_rgb(color))
        run2.font.name = 'SimSun-ExtB'
    if fill:
        tcPr = cell._tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', fill)
    set_borders(cell)
    cell.margin_left = Pt(1); cell.margin_right = Pt(1)
    cell.margin_top = Pt(0.5); cell.margin_bottom = Pt(0.5)

GREEN = 'DAF2D0'; BLUE = 'DAE9F8'; YELLOW = 'FFFF00'; GREY = 'F2F2F2'

def process(input_pptx, output_pptx):
    prs = Presentation(input_pptx)
    slide5 = prs.slides[4]

    # Remove existing images
    to_remove = [s for s in slide5.shapes if s.shape_type == 13]
    for s in to_remove:
        s._element.getparent().remove(s._element)
    print(f'Removed {len(to_remove)} images from slide 5')

    # --- Table 1: Weekly Overview (8x3) ---
    t1 = slide5.shapes.add_table(8, 3, Emu(287718), Emu(1735390), Emu(3397250), Emu(2400000)).table
    t1.columns[0].width = Emu(883578); t1.columns[1].width = Emu(1328328); t1.columns[2].width = Emu(1506762)

    overview = [
        [("Overview 03/07/26", "概览 2026年7月3日", 7, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("SN1003", "", 7, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("SN1004", "", 7, True, GREEN, '000000', PP_ALIGN.CENTER)],
        [("WEEKLY", "周度", 7, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("", "", 7, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("", "", 7, True, GREEN, '000000', PP_ALIGN.CENTER)],
        [("Monday", "周一", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("200FH & Yearly Maintenance", "200小时及年检维护", 6, False, None, '000000', PP_ALIGN.LEFT),
         ("1*Flight Test SW Content Freeze\nRemaining Items\n(No Flight due to Weather)",
          "1*飞行测试：软件内容冻结剩余项目\n（因天气原因未飞行）", 6, False, None, '000000', PP_ALIGN.LEFT)],
        [("Tuesday", "周二", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("200FH & Yearly Maintenance", "200小时及年检维护", 6, False, None, '000000', PP_ALIGN.LEFT),
         ("10FH inspection + FTI TRS", "10小时检查 + FTI TRS", 6, False, None, '000000', PP_ALIGN.LEFT)],
        [("Wednesday", "周三", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("200FH & Yearly Maintenance", "200小时及年检维护", 6, False, None, '000000', PP_ALIGN.LEFT),
         ("3*Flight Test Avionics Cert.\nfor Content Freeze",
          "3*飞行测试：航电认证\n支持内容冻结", 6, False, None, '000000', PP_ALIGN.LEFT)],
        [("Thursday", "周四", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("200FH & Yearly Maintenance", "200小时及年检维护", 6, False, None, '000000', PP_ALIGN.LEFT),
         ("1*Flight Test Avionics Cert.\nfor Content Freeze",
          "1*飞行测试：航电认证\n支持内容冻结", 6, False, None, '000000', PP_ALIGN.LEFT)],
        [("Friday", "周五", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("200FH & Yearly Maintenance", "200小时及年检维护", 6, False, None, '000000', PP_ALIGN.LEFT),
         ("1*Flight Test SW Content Freeze", "1*飞行测试：软件内容冻结", 6, False, None, '000000', PP_ALIGN.LEFT)],
        [("WEEKLY", "周度", 6, True, GREY, '000000', PP_ALIGN.LEFT),
         ("BLOCK: 0:00\nFLIGHT: 0:00", "轮挡: 0:00 / 飞行: 0:00", 6, False, GREY, '000000', PP_ALIGN.CENTER),
         ("BLOCK: 7:29\nFLIGHT: 6:16", "轮挡: 7:29 / 飞行: 6:16", 6, False, GREY, '000000', PP_ALIGN.CENTER)],
    ]
    for ri, row in enumerate(overview):
        for ci, (en, cn, sz, b, fill, color, align) in enumerate(row):
            set_cell(t1.cell(ri, ci), en, cn, Pt(sz), b, color, fill, align)
    print('Table 1 (Weekly Overview): 8x3')

    # --- Table 2: Flight Hours & Flights (5x9) ---
    t2 = slide5.shapes.add_table(5, 9, Emu(3727224), Emu(1729740), Emu(4140200), Emu(1682590)).table
    for ci in range(9):
        t2.columns[ci].width = Emu(296177 if ci == 0 else 427811)

    fh_rows = [
        [("Flight hours", "飞行小时"),("",""),("",""),("",""),("",""),
         ("Flights", "次数"),("",""),("",""),("","")],
        [("SN",""),("2024",""),("2025",""),("2026",""),("Total","累计"),
         ("2024",""),("2025",""),("2026",""),("Total","累计")],
        [("1002",""),("49:32",""),("0:00",""),("0",""),("1237:09",""),
         ("48",""),("0",""),("0",""),("48","")],
        [("1003",""),("16:31",""),("64:32",""),("0",""),("198:27",""),
         ("16",""),("55",""),("0",""),("71","")],
        [("1004",""),("27:52",""),("80:52",""),("13:44",""),("122:28",""),
         ("32",""),("89",""),("13",""),("134","")],
    ]
    for ri, row in enumerate(fh_rows):
        for ci, (en, cn) in enumerate(row):
            set_cell(t2.cell(ri, ci), en, cn, Pt(6), ri <= 1, '000000', BLUE if ri <= 1 else None)
    # Merge header cells
    t2.cell(0,0).merge(t2.cell(0,4)); t2.cell(0,5).merge(t2.cell(0,8))
    print('Table 2 (Flight Hours): 5x9')

    # --- Table 3: Progress Summary (5x4) ---
    t3 = slide5.shapes.add_table(5, 4, Emu(7847331), Emu(1729740), Emu(4026898), Emu(1682590)).table
    for ci in range(4):
        t3.columns[ci].width = Emu(4026898 // 4)

    prog_data = [
        [("Total scheduled\nflight hours", "计划飞行总小时数"), ("Hours flown so far", "已完成飞行小时数"),
         ("Percentage\ncompleted", "完成百分比"), ("Hours to go\n(some need to\nbe repeated)", "剩余小时数\n（部分需重复）")],
        [("500:00",""),("320:55",""),("64%",""),("179:05","")],
        [("","")]*4,
        [("Total task items\n(FRQ's) scheduled", "计划任务项(FRQ)总数"), ("FRQ's so far\ncompleted", "目前已完成"),
         ("Percentage FRQ's\ncompleted", "FRQ完成百分比"), ("FRQ's started", "已开始")],
        [("202",""),("100",""),("50%",""),("120","")],
    ]
    for ri, row in enumerate(prog_data):
        for ci, (en, cn) in enumerate(row):
            is_h = (ri in [0, 3])
            set_cell(t3.cell(ri, ci), en, cn, Pt(6), is_h, '000000', BLUE if is_h else None)
    print('Table 3 (Progress Summary): 5x4')

    # --- Table 4: Block/Flight Time (3x3) ---
    t4 = slide5.shapes.add_table(3, 3, Emu(287718), Emu(4200000), Emu(3397250), Emu(585000)).table
    t4.columns[0].width = Emu(1524232); t4.columns[1].width = Emu(1114527); t4.columns[2].width = Emu(1070343)

    block_data = [
        [("WEEKLY", "周度", 6, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("SN1003", "", 6, True, GREEN, '000000', PP_ALIGN.CENTER),
         ("SN1004", "", 6, True, GREEN, '000000', PP_ALIGN.CENTER)],
        [("BLOCK TIME\n(chock to chock)", "轮挡时间\n（轮挡到轮挡）", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("0:00", "", 6, False, YELLOW, '000000', PP_ALIGN.CENTER),
         ("7:29", "", 6, False, YELLOW, '000000', PP_ALIGN.CENTER)],
        [("FLIGHT TIME\n(liftoff to touchdown)", "飞行时间\n（起飞到着陆）", 6, True, None, '000000', PP_ALIGN.LEFT),
         ("0:00", "", 6, False, YELLOW, '000000', PP_ALIGN.CENTER),
         ("6:16", "", 6, False, YELLOW, '000000', PP_ALIGN.CENTER)],
    ]
    for ri, row in enumerate(block_data):
        for ci, (en, cn, sz, b, fill, color, align) in enumerate(row):
            set_cell(t4.cell(ri, ci), en, cn, Pt(sz), b, color, fill, align)
    print('Table 4 (Block/Flight Time): 3x3')

    # Save
    prs.save(output_pptx)
    print(f'Saved: {output_pptx}')

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: python rebuild_flight_test_progress.py <input.pptx> <output.pptx>')
        sys.exit(1)
    process(sys.argv[1], sys.argv[2])
